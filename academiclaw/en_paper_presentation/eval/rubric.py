"""
Rubric for generate/ppt/junjieyu-query3
Task: Generate a detailed analysis PPT based on multiple academic papers (GPT-1/2/3/4)

Total: 100 points

Scoring Dimensions:
  1. File Delivery          (15 pts) — deterministic file existence & validity checks
  2. PPT Structure & Design (25 pts) — deterministic python-pptx inspection
  3. Code Quality           (10 pts) — AST parsing + keyword analysis of generate_ppt.py
  4. Summary Report Quality  (25 pts) — LLM-as-Judge on papers_summary.md
  5. PPT Content Quality     (25 pts) — LLM-as-Judge on extracted PPT text
"""

import os
import re
import ast
import json
from typing import Tuple, Dict, Any, Optional

try:
    import openai
except ImportError:
    openai = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# ── Papers that must be covered ──────────────────────────────────────────────
KNOWN_PAPERS = [
    "Improving Language Understanding by Generative Pre-Training",  # GPT-1
    "Language Models are Unsupervised Multitask Learners",          # GPT-2
    "Language Models are Few-Shot Learners",                        # GPT-3
    "GPT-4 Technical Report",                                      # GPT-4
]
PAPER_KEYWORDS = [
    "gpt-1", "gpt-2", "gpt-3", "gpt-4", "gpt4",
    "pre-training", "few-shot", "multitask", "technical report",
]


# ═══════════════════════════════════════════════════════════════════════════════
# Helper: Environment & LLM
# ═══════════════════════════════════════════════════════════════════════════════

def _load_env(answer_dir: str) -> dict:
    """Load .env values from answer_dir and query root."""
    values: Dict[str, str] = {}
    for env_dir in [answer_dir, os.path.join(os.path.dirname(__file__), "..")]:
        env_path = os.path.join(env_dir, ".env")
        if os.path.exists(env_path):
            try:
                with open(env_path, "r") as f:
                    for line in f:
                        line = line.strip()
                        if not line or line.startswith("#") or "=" not in line:
                            continue
                        k, v = line.split("=", 1)
                        if k.strip() not in values:
                            values[k.strip()] = v.strip().strip("'\"")
            except Exception:
                pass
    return values


def _get_text_eval_config(answer_dir: str) -> dict:
    env = _load_env(answer_dir)
    def g(key, default=""):
        return os.environ.get(key) or env.get(key) or default
    return {
        "api_key": g("EVAL_TEXT_API_KEY", g("ANTHROPIC_API_KEY")),
        "api_base": g("EVAL_TEXT_API_BASE_URL", g("ANTHROPIC_BASE_URL")),
        "model": g("EVAL_TEXT_MODEL", "openai/gpt-5.2"),
    }


def _call_llm_judge(prompt: str, config: dict) -> str:
    if not openai or not config.get("api_key"):
        return ""
    try:
        base = config["api_base"].rstrip("/")
        if not base.endswith("/v1"):
            base += "/v1"
        client = openai.OpenAI(api_key=config["api_key"], base_url=base)
        resp = client.chat.completions.create(
            model=config["model"],
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2048,
            temperature=0,
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        print(f"[RUBRIC] LLM Judge call failed: {e}")
        return ""


def _parse_json_from_llm(text: str) -> Optional[dict]:
    if not text:
        return None
    m = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
    raw = m.group(1) if m else text
    try:
        return json.loads(raw)
    except Exception:
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# Helper: File I/O
# ═══════════════════════════════════════════════════════════════════════════════

def _find_file(answer_dir: str, primary: str, ext: str) -> Optional[str]:
    """Return path of primary file if it exists, else the first file matching ext."""
    p = os.path.join(answer_dir, primary)
    if os.path.isfile(p):
        return p
    try:
        for f in os.listdir(answer_dir):
            if f.lower().endswith(ext):
                return os.path.join(answer_dir, f)
    except Exception:
        pass
    return None


def _read_text(path: str, max_chars: int = 0) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return content[:max_chars] if max_chars else content
    except Exception:
        return ""


def _load_pptx(path: str):
    if Presentation is None or not path:
        return None
    try:
        return Presentation(path)
    except Exception:
        return None


def _extract_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            parts.append(shape.text)
    return " ".join(parts)


def _extract_all_ppt_text(prs) -> str:
    sections = []
    for i, slide in enumerate(prs.slides):
        txt = _extract_slide_text(slide)
        if txt.strip():
            sections.append(f"=== Slide {i + 1} ===\n{txt}")
    return "\n\n".join(sections)


# ═══════════════════════════════════════════════════════════════════════════════
# Dimension 1: File Delivery (15 pts)
# ═══════════════════════════════════════════════════════════════════════════════

def _eval_files(answer_dir: str) -> Tuple[int, dict]:
    """
    papers_summary.md   — 4 pts
    papers_analysis.pptx — 4 pts
    ppt_structure.json  — 3 pts
    generate_ppt.py     — 4 pts
    """
    score = 0
    details: Dict[str, str] = {}

    try:
        all_files = os.listdir(answer_dir)
    except Exception:
        return 0, {"error": "cannot list answer_dir"}

    # 1.1 papers_summary.md (4 pts)
    md_path = _find_file(answer_dir, "papers_summary.md", ".md")
    if md_path:
        content = _read_text(md_path)
        if len(content) >= 500:
            score += 4
            details["papers_summary.md"] = f"4/4 — found ({len(content)} chars)"
        elif len(content) >= 100:
            score += 2
            details["papers_summary.md"] = f"2/4 — found but short ({len(content)} chars)"
        else:
            score += 1
            details["papers_summary.md"] = f"1/4 — found but nearly empty ({len(content)} chars)"
    else:
        details["papers_summary.md"] = "0/4 — missing"

    # 1.2 papers_analysis.pptx (4 pts)
    pptx_path = _find_file(answer_dir, "papers_analysis.pptx", ".pptx")
    if pptx_path:
        fsize = os.path.getsize(pptx_path)
        prs = _load_pptx(pptx_path)
        if prs is not None and fsize >= 10240:
            score += 4
            details["papers_analysis.pptx"] = f"4/4 — valid PPTX ({fsize / 1024:.0f} KB)"
        elif prs is not None:
            score += 3
            details["papers_analysis.pptx"] = f"3/4 — valid but small ({fsize} bytes)"
        elif fsize >= 10240:
            score += 2
            details["papers_analysis.pptx"] = f"2/4 — file exists ({fsize / 1024:.0f} KB) but cannot parse"
        else:
            score += 1
            details["papers_analysis.pptx"] = f"1/4 — file exists but too small and cannot parse"
    else:
        details["papers_analysis.pptx"] = "0/4 — missing"

    # 1.3 ppt_structure.json (3 pts)
    json_path = _find_file(answer_dir, "ppt_structure.json", ".json")
    if json_path:
        try:
            with open(json_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            if isinstance(data, dict) and "papers" in data:
                score += 3
                details["ppt_structure.json"] = "3/3 — valid JSON with 'papers' key"
            elif isinstance(data, dict):
                score += 2
                details["ppt_structure.json"] = "2/3 — valid JSON but missing 'papers' key"
            else:
                score += 1
                details["ppt_structure.json"] = "1/3 — valid JSON but unexpected structure"
        except Exception as e:
            score += 1
            details["ppt_structure.json"] = f"1/3 — file exists but invalid JSON: {str(e)[:50]}"
    else:
        details["ppt_structure.json"] = "0/3 — missing"

    # 1.4 generate_ppt.py (4 pts)
    py_path = _find_file(answer_dir, "generate_ppt.py", ".py")
    if py_path:
        code = _read_text(py_path)
        try:
            ast.parse(code)
            if len(code) >= 500:
                score += 4
                details["generate_ppt.py"] = f"4/4 — syntax OK, {len(code)} chars"
            elif len(code) >= 100:
                score += 3
                details["generate_ppt.py"] = f"3/4 — syntax OK but short ({len(code)} chars)"
            else:
                score += 2
                details["generate_ppt.py"] = f"2/4 — syntax OK but very short ({len(code)} chars)"
        except SyntaxError as e:
            score += 1
            details["generate_ppt.py"] = f"1/4 — syntax error: {str(e)[:50]}"
    else:
        details["generate_ppt.py"] = "0/4 — missing"

    return score, details


# ═══════════════════════════════════════════════════════════════════════════════
# Dimension 2: PPT Structure & Design (25 pts)
# ═══════════════════════════════════════════════════════════════════════════════

def _eval_ppt_structure(answer_dir: str) -> Tuple[int, dict]:
    """
    2.1  Slide count            (5 pts) — 4 papers × ~9 sections + cover/TOC/overview/compare/end ≈ 40+
    2.2  Cover page             (3 pts) — title + date
    2.3  Table of contents      (3 pts)
    2.4  Separator pages        (4 pts) — one per paper
    2.5  Ending page            (2 pts) — thank you / Q&A
    2.6  16:9 dimensions        (3 pts)
    2.7  Font sizes             (3 pts) — title ≥28pt, body ≥10pt
    2.8  Content density        (2 pts) — not overcrowded
    """
    pptx_path = _find_file(answer_dir, "papers_analysis.pptx", ".pptx")
    prs = _load_pptx(pptx_path)
    if prs is None:
        return 0, {"error": "0/25 — cannot open PPT (python-pptx missing or file invalid)"}

    score = 0
    details: Dict[str, str] = {}
    slides = list(prs.slides)
    total_slides = len(slides)

    # 2.1 Slide count (5 pts)
    if total_slides >= 36:
        score += 5
        details["slide_count"] = f"5/5 — {total_slides} slides (>=36)"
    elif total_slides >= 25:
        score += 3
        details["slide_count"] = f"3/5 — {total_slides} slides (>=25 but <36)"
    elif total_slides >= 10:
        score += 1
        details["slide_count"] = f"1/5 — {total_slides} slides (too few for 4 papers)"
    else:
        details["slide_count"] = f"0/5 — only {total_slides} slides"

    # 2.2 Cover page (3 pts)
    if total_slides > 0:
        cover_txt = _extract_slide_text(slides[0]).lower()
        has_title = len(cover_txt) > 10
        has_date = any(kw in cover_txt for kw in ["2025", "2026", "日期", "date", "生成"])
        if has_title and has_date:
            score += 3
            details["cover_page"] = "3/3 — title + date present"
        elif has_title:
            score += 2
            details["cover_page"] = "2/3 — title present, no date"
        else:
            score += 0
            details["cover_page"] = "0/3 — cover page incomplete"
    else:
        details["cover_page"] = "0/3 — no slides"

    # 2.3 Table of contents (3 pts)
    toc_found = False
    for s in slides[:6]:
        txt = _extract_slide_text(s).lower()
        if any(kw in txt for kw in ["目录", "contents", "outline", "table of contents", "agenda"]):
            toc_found = True
            break
    details["toc_page"] = "3/3 — found" if toc_found else "0/3 — not found"
    if toc_found:
        score += 3

    # 2.4 Separator pages for each paper (4 pts)
    separator_count = 0
    for s in slides:
        txt = _extract_slide_text(s).lower()
        shape_count = len([sh for sh in s.shapes if hasattr(sh, "text") and sh.text.strip()])
        # Separator pages: short text, few shapes, mentions a paper
        if len(txt) < 250 and shape_count <= 5:
            if any(kw in txt for kw in [
                "et al", "radford", "brown", "openai",
                "gpt-1", "gpt-2", "gpt-3", "gpt-4", "gpt1", "gpt2", "gpt3", "gpt4",
                "few-shot", "multitask", "pre-training", "technical report",
                "论文", "paper",
            ]):
                separator_count += 1
    if separator_count >= 4:
        score += 4
        details["separator_pages"] = f"4/4 — {separator_count} separator pages found"
    elif separator_count >= 2:
        score += 2
        details["separator_pages"] = f"2/4 — only {separator_count} separator pages"
    elif separator_count >= 1:
        score += 1
        details["separator_pages"] = f"1/4 — only {separator_count} separator page"
    else:
        details["separator_pages"] = "0/4 — no separator pages detected"

    # 2.5 Ending page (2 pts)
    if total_slides > 0:
        last_txt = _extract_slide_text(slides[-1]).lower()
        if any(kw in last_txt for kw in ["谢谢", "thank", "q&a", "question", "讨论", "q & a"]):
            score += 2
            details["ending_page"] = "2/2 — found"
        else:
            details["ending_page"] = "0/2 — not found"
    else:
        details["ending_page"] = "0/2 — no slides"

    # 2.6 Slide dimensions — 16:9 (3 pts)
    w_in = prs.slide_width.inches
    h_in = prs.slide_height.inches
    ratio = w_in / h_in if h_in else 0
    if 1.7 <= ratio <= 1.85:
        score += 3
        details["dimensions"] = f"3/3 — 16:9 ({w_in:.2f} x {h_in:.2f} in)"
    elif 1.3 <= ratio <= 1.4:
        score += 1
        details["dimensions"] = f"1/3 — 4:3 ({w_in:.2f} x {h_in:.2f} in), expected 16:9"
    else:
        details["dimensions"] = f"0/3 — non-standard ratio {ratio:.2f}"

    # 2.7 Font sizes (3 pts)
    font_sizes = []
    for s in slides[:20]:
        for shape in s.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)
    if font_sizes:
        mx = max(font_sizes)
        mn = min(font_sizes)
        if mx >= 28 and mn >= 10:
            score += 3
            details["font_sizes"] = f"3/3 — range {mn:.0f}-{mx:.0f} pt"
        elif mx >= 24:
            score += 2
            details["font_sizes"] = f"2/3 — range {mn:.0f}-{mx:.0f} pt (slightly small titles)"
        else:
            score += 1
            details["font_sizes"] = f"1/3 — range {mn:.0f}-{mx:.0f} pt (too small)"
    else:
        score += 1
        details["font_sizes"] = "1/3 — unable to detect explicit font sizes"

    # 2.8 Content density (2 pts)
    lengths = [len(_extract_slide_text(s)) for s in slides]
    if lengths:
        max_len = max(lengths)
        avg_len = sum(lengths) / len(lengths)
        if max_len <= 1000 and avg_len <= 500:
            score += 2
            details["content_density"] = f"2/2 — avg {avg_len:.0f}, max {max_len}"
        elif max_len <= 1500:
            score += 1
            details["content_density"] = f"1/2 — some slides dense (max {max_len} chars)"
        else:
            details["content_density"] = f"0/2 — overcrowded (max {max_len} chars)"
    else:
        details["content_density"] = "0/2 — no content"

    return score, details


# ═══════════════════════════════════════════════════════════════════════════════
# Dimension 3: Code Quality (10 pts)
# ═══════════════════════════════════════════════════════════════════════════════

def _eval_code(answer_dir: str) -> Tuple[int, dict]:
    """
    3.1  Syntax correct         (3 pts)
    3.2  Uses python-pptx       (3 pts)
    3.3  References papers      (2 pts)
    3.4  Saves output files     (2 pts)
    """
    py_path = _find_file(answer_dir, "generate_ppt.py", ".py")
    if not py_path:
        return 0, {"error": "0/10 — generate_ppt.py not found"}

    code = _read_text(py_path)
    if not code:
        return 0, {"error": "0/10 — file is empty"}

    score = 0
    details: Dict[str, str] = {}

    # 3.1 Syntax (3 pts)
    try:
        ast.parse(code)
        score += 3
        details["syntax"] = "3/3 — no syntax errors"
    except SyntaxError as e:
        details["syntax"] = f"0/3 — {str(e)[:80]}"
        return score, details

    code_lower = code.lower()

    # 3.2 Uses python-pptx (3 pts)
    if "from pptx" in code_lower or "import pptx" in code_lower:
        score += 3
        details["uses_pptx"] = "3/3 — imports python-pptx"
    elif "pptx" in code_lower:
        score += 1
        details["uses_pptx"] = "1/3 — references pptx but no clear import"
    else:
        details["uses_pptx"] = "0/3 — does not use python-pptx"

    # 3.3 Paper references (2 pts)
    paper_kw_count = sum(1 for kw in [
        "gpt-1", "gpt-2", "gpt-3", "gpt-4", "gpt1", "gpt2", "gpt3", "gpt4",
        "few-shot", "multitask", "pre-training", "technical report",
    ] if kw in code_lower)
    if paper_kw_count >= 3:
        score += 2
        details["paper_refs"] = f"2/2 — {paper_kw_count} paper keywords found"
    elif paper_kw_count >= 1:
        score += 1
        details["paper_refs"] = f"1/2 — {paper_kw_count} paper keyword(s) found"
    else:
        # Still give 1pt if code is substantial (dynamic reading)
        if len(code) > 3000:
            score += 1
            details["paper_refs"] = "1/2 — no explicit keywords but code is substantial"
        else:
            details["paper_refs"] = "0/2 — no paper references"

    # 3.4 Saves output files (2 pts)
    saves = 0
    if re.search(r'\.save\s*\(', code) or "papers_analysis.pptx" in code:
        saves += 1
    if "ppt_structure.json" in code or "papers_summary" in code_lower:
        saves += 1
    score += min(2, saves)
    details["output_files"] = f"{min(2, saves)}/2 — {saves} output reference(s)"

    return score, details


# ═══════════════════════════════════════════════════════════════════════════════
# Dimension 4: Summary Report Quality — LLM Judge (25 pts)
# ═══════════════════════════════════════════════════════════════════════════════

_SUMMARY_PROMPT = """\
You are a strict academic reviewer evaluating a paper analysis report.

The report should cover 4 papers about the GPT model series:
1. GPT-1: "Improving Language Understanding by Generative Pre-Training" (Radford et al.)
2. GPT-2: "Language Models are Unsupervised Multitask Learners" (Radford et al.)
3. GPT-3: "Language Models are Few-Shot Learners" (Brown et al., 2020)
4. GPT-4: "GPT-4 Technical Report" (OpenAI, 2024)

The task requires the report to contain:
- Overview section with research field summary, cross-paper connections, and trends
- Per-paper detailed analysis covering: basic info, background & motivation, core contributions (3-5 points), methods, experiments, pros/cons, personal insights
- Cross-paper comparison table
- Summary and future outlook

## Evaluation Criteria (integer scores only):

1. **paper_coverage** (0-8): Does the report cover ALL 4 papers with accurate basic info (title, authors, venue, year)?
   - 8: all 4 papers thoroughly covered with correct metadata
   - 5-7: 3-4 papers covered but some metadata missing
   - 2-4: only 1-2 papers covered
   - 0-1: no meaningful paper content

2. **analysis_depth** (0-7): Does each paper analysis include: background/motivation, core contributions (3-5 points), detailed methods, experiments with specific results, pros/cons, and personal insights?
   - 7: all sections present for most papers with specific technical details
   - 4-6: most sections present but lacks specifics
   - 1-3: shallow summary only
   - 0: no real analysis

3. **cross_paper_analysis** (0-5): Is there a clear cross-paper comparison, evolution narrative (GPT-1→2→3→4), and comparison table?
   - 5: excellent comparison with table and evolution analysis
   - 3-4: some comparison but incomplete
   - 1-2: minimal cross-paper discussion
   - 0: no cross-paper analysis

4. **writing_quality** (0-5): Is the report well-structured, clearly written, and comprehensive?
   - 5: professional quality, clear structure, complete coverage
   - 3-4: acceptable but some structural issues
   - 1-2: poor structure or missing major sections
   - 0: incoherent or extremely short

## Content to evaluate:
{content}

Return ONLY a JSON object (no markdown fences):
{{"paper_coverage": {{"score": 0, "reason": ""}}, "analysis_depth": {{"score": 0, "reason": ""}}, "cross_paper_analysis": {{"score": 0, "reason": ""}}, "writing_quality": {{"score": 0, "reason": ""}}, "total": 0}}
"""


def _eval_summary(answer_dir: str) -> Tuple[int, dict]:
    md_path = _find_file(answer_dir, "papers_summary.md", ".md")
    if not md_path:
        return 0, {"error": "0/25 — papers_summary.md not found"}

    content = _read_text(md_path)
    if len(content) < 100:
        return 0, {"error": "0/25 — summary too short to evaluate"}

    # Deterministic fallback scoring (max 7/25)
    det_score = 0
    det_details: Dict[str, str] = {}

    content_lower = content.lower()
    # Check paper keyword coverage
    kw_hits = sum(1 for kw in PAPER_KEYWORDS if kw in content_lower)
    if kw_hits >= 5:
        det_score += 3
        det_details["paper_coverage_det"] = f"3 — {kw_hits} paper keywords found"
    elif kw_hits >= 2:
        det_score += 1
        det_details["paper_coverage_det"] = f"1 — only {kw_hits} paper keywords"
    else:
        det_details["paper_coverage_det"] = f"0 — only {kw_hits} paper keywords"

    # Check structural sections
    has_overview = any(kw in content_lower for kw in ["综述", "overview", "概述", "研究领域"])
    has_comparison = any(kw in content_lower for kw in ["对比", "comparison", "比较", "对比分析"])
    has_future = any(kw in content_lower for kw in ["展望", "future", "未来", "outlook"])
    if has_overview:
        det_score += 1
    if has_comparison:
        det_score += 1
    if has_future:
        det_score += 1
    det_details["structure_det"] = (
        f"overview={has_overview}, comparison={has_comparison}, future={has_future}"
    )

    # Length check
    if len(content) >= 5000:
        det_score += 1
        det_details["length_det"] = f"1 — substantial ({len(content)} chars)"
    else:
        det_details["length_det"] = f"0 — short ({len(content)} chars)"

    # LLM evaluation
    config = _get_text_eval_config(answer_dir)
    prompt = _SUMMARY_PROMPT.format(content=content[:10000])
    raw = _call_llm_judge(prompt, config)
    result = _parse_json_from_llm(raw)

    if result:
        dims = {
            "paper_coverage": 8,
            "analysis_depth": 7,
            "cross_paper_analysis": 5,
            "writing_quality": 5,
        }
        llm_score = 0
        details: Dict[str, str] = {}
        for dim, cap in dims.items():
            entry = result.get(dim, {})
            s = max(0, min(cap, int(entry.get("score", 0))))
            llm_score += s
            details[dim] = f"{s}/{cap} — {entry.get('reason', '')[:120]}"
        details["llm_total"] = f"{llm_score}/25"
        return llm_score, details
    else:
        det_details["note"] = "LLM unavailable, using deterministic fallback (max 7/25)"
        return det_score, det_details


# ═══════════════════════════════════════════════════════════════════════════════
# Dimension 5: PPT Content Quality — LLM Judge (25 pts)
# ═══════════════════════════════════════════════════════════════════════════════

_PPT_CONTENT_PROMPT = """\
You are a strict reviewer evaluating an academic paper analysis PPT.

The PPT should cover 4 papers about the GPT series:
1. GPT-1: "Improving Language Understanding by Generative Pre-Training"
2. GPT-2: "Language Models are Unsupervised Multitask Learners"
3. GPT-3: "Language Models are Few-Shot Learners"
4. GPT-4: "GPT-4 Technical Report"

For EACH paper, the PPT should include slides on:
- Paper basic info (title, authors, venue, year)
- Background & motivation
- Core contributions (3-5 points)
- Methods / model architecture
- Experiment setup
- Experiment results (ideally with tables)
- Analysis & discussion
- Limitations & future work
- Summary

Additionally there should be:
- A cover slide, table of contents, overview/comparison slides, and ending slide

## Evaluation Criteria (integer scores only):

1. **content_accuracy** (0-10): Does the PPT accurately reflect paper content? Are key technical details (model sizes, training data, benchmark results) correct? No hallucinations?
   - 8-10: highly accurate, specific technical details present
   - 5-7: mostly accurate but lacks specifics
   - 2-4: significant inaccuracies or very shallow
   - 0-1: no meaningful content

2. **paper_completeness** (0-8): Are ALL 4 papers covered with the required sections (background, contributions, methods, experiments, analysis, limitations, summary)?
   - 7-8: all 4 papers with most required sections
   - 4-6: 3-4 papers but some sections missing
   - 1-3: only 1-2 papers or very incomplete
   - 0: no paper-specific content

3. **presentation_quality** (0-7): Is the content well-organized? Logical flow between slides? Good use of bullet points and tables? Not too dense?
   - 6-7: professional quality, good organization
   - 3-5: acceptable but some organizational issues
   - 1-2: poor organization or walls of text
   - 0: unusable

## PPT Content (extracted text from slides):
{ppt_text}

Return ONLY a JSON object (no markdown fences):
{{"content_accuracy": {{"score": 0, "reason": ""}}, "paper_completeness": {{"score": 0, "reason": ""}}, "presentation_quality": {{"score": 0, "reason": ""}}, "total": 0}}
"""


def _eval_ppt_content(answer_dir: str) -> Tuple[int, dict]:
    pptx_path = _find_file(answer_dir, "papers_analysis.pptx", ".pptx")
    prs = _load_pptx(pptx_path)
    if prs is None:
        return 0, {"error": "0/25 — cannot open PPT"}

    ppt_text = _extract_all_ppt_text(prs)
    if len(ppt_text) < 100:
        return 0, {"error": "0/25 — PPT has almost no text"}

    # Deterministic fallback (max 7/25)
    det_score = 0
    det_details: Dict[str, str] = {}
    text_lower = ppt_text.lower()

    # Check paper mention coverage
    papers_found = 0
    for paper_title in KNOWN_PAPERS:
        if paper_title.lower()[:25] in text_lower:
            papers_found += 1
    if papers_found >= 4:
        det_score += 3
    elif papers_found >= 2:
        det_score += 1
    det_details["papers_in_ppt"] = f"{papers_found}/4 papers mentioned"

    # Check key sections
    section_kws = [
        "背景", "贡献", "方法", "实验", "结果", "局限", "总结",
        "background", "contribution", "method", "experiment",
        "result", "limitation", "architecture",
    ]
    sections_found = sum(1 for kw in section_kws if kw in text_lower)
    if sections_found >= 6:
        det_score += 3
    elif sections_found >= 3:
        det_score += 1
    det_details["sections_found"] = f"{sections_found} section keywords"

    # Length bonus
    if len(ppt_text) >= 3000:
        det_score += 1
        det_details["text_length"] = f"1 — substantial text ({len(ppt_text)} chars)"

    # LLM evaluation
    config = _get_text_eval_config(answer_dir)
    prompt = _PPT_CONTENT_PROMPT.format(ppt_text=ppt_text[:10000])
    raw = _call_llm_judge(prompt, config)
    result = _parse_json_from_llm(raw)

    if result:
        dims = {
            "content_accuracy": 10,
            "paper_completeness": 8,
            "presentation_quality": 7,
        }
        llm_score = 0
        details: Dict[str, str] = {}
        for dim, cap in dims.items():
            entry = result.get(dim, {})
            s = max(0, min(cap, int(entry.get("score", 0))))
            llm_score += s
            details[dim] = f"{s}/{cap} — {entry.get('reason', '')[:120]}"
        details["llm_total"] = f"{llm_score}/25"
        return llm_score, details
    else:
        det_details["note"] = "LLM unavailable, using deterministic fallback (max 7/25)"
        return det_score, det_details


# ═══════════════════════════════════════════════════════════════════════════════
# Main entry points
# ═══════════════════════════════════════════════════════════════════════════════

def evaluate(answer_dir: str) -> Tuple[int, Dict[str, Any]]:
    """
    Evaluate agent output for the academic paper analysis PPT task.

    Args:
        answer_dir: absolute path to the agent output directory

    Returns:
        (score, report) where score is 0-100 and report is a detailed dict
    """
    s1, r1 = _eval_files(answer_dir)
    s2, r2 = _eval_ppt_structure(answer_dir)
    s3, r3 = _eval_code(answer_dir)
    s4, r4 = _eval_summary(answer_dir)
    s5, r5 = _eval_ppt_content(answer_dir)

    total = s1 + s2 + s3 + s4 + s5

    report = {
        "total": total,
        "breakdown": {
            "1_file_delivery (15)": {"score": s1, "details": r1},
            "2_ppt_structure_design (25)": {"score": s2, "details": r2},
            "3_code_quality (10)": {"score": s3, "details": r3},
            "4_summary_report_quality (25)": {"score": s4, "details": r4},
            "5_ppt_content_quality (25)": {"score": s5, "details": r5},
        },
    }

    if total >= 90:
        report["comment"] = "Excellent — thorough analysis, professional PPT, strong insights."
    elif total >= 75:
        report["comment"] = "Good — task mostly completed with minor gaps."
    elif total >= 60:
        report["comment"] = "Acceptable — core deliverables present but notable shortcomings."
    elif total >= 40:
        report["comment"] = "Partial — significant issues in quality or missing deliverables."
    else:
        report["comment"] = "Insufficient — major deliverables missing or fundamentally flawed."

    return total, report


def print_report(score: int, report: Dict[str, Any]) -> None:
    """Print a formatted evaluation report."""
    print("=" * 70)
    print("Evaluation Report — Academic Paper Analysis PPT")
    print("Task: generate/ppt/junjieyu-query3")
    print("=" * 70)
    print(f"\nTotal Score: {score}/100\n")

    for section_key, section_data in report.get("breakdown", {}).items():
        sec_score = section_data.get("score", 0)
        print(f"--- {section_key}: {sec_score} pts ---")
        details = section_data.get("details", {})
        if isinstance(details, dict):
            for k, v in details.items():
                v_str = str(v)
                if len(v_str) > 130:
                    v_str = v_str[:130] + "..."
                print(f"  {k}: {v_str}")
        else:
            print(f"  {details}")
        print()

    comment = report.get("comment", "")
    if comment:
        print(f"Comment: {comment}")
    print("=" * 70)


if __name__ == "__main__":
    import sys
    test_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(__file__), "..", "gpt-5", "attempt_1")
    test_dir = os.path.abspath(test_dir)
    if os.path.exists(test_dir):
        print(f"Evaluating: {test_dir}\n")
        s, r = evaluate(test_dir)
        print_report(s, r)
    else:
        print(f"Directory not found: {test_dir}")
    sys.exit(0)
